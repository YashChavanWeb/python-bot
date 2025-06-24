from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from typing import List
import os
import google.generativeai as genai
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import io
import re
import sounddevice as sd
from scipy.io.wavfile import write
import json
from datetime import datetime
import speech_recognition as sr
from pydantic import BaseModel  # Import BaseModel for pydantic


# ==== Gemini Config ====
API_KEY = "AIzaSyA58up6mb0EppG3dI0lT2WYct4Om9aEQKw"  # Your actual Gemini API key
genai.configure(api_key=API_KEY)
model = genai.GenerativeModel(
    "gemini-2.0-flash",
    system_instruction="Your responses should be concise, short, direct, and educational, answering only educational queries with a 120 to 200 word limit.",
)

# ==== FastAPI App ====
app = FastAPI()

# CORS configuration
origins = [
    "http://localhost",
    "http://localhost:8000",
    "http://localhost:5173",  # React development server
    # You might want to add other specific origins for production
    "*",  # This is for development, be more restrictive in production
]

# Enable CORS for frontend calls
app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Session state for Quizzy Bot
chat_history = []
document_context = ""

# Constants for Audio Recording
DURATION = 10  # seconds to record
SAMPLE_RATE = 44100
AUDIO_DIR = "audio_files"
OUTPUT_DIR = "outputs"
AUDIO_FILENAME = "audio.wav"
JSON_FILENAME = "transcription.json"

# Create directories if they don't exist
os.makedirs(AUDIO_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

audio_path = os.path.join(AUDIO_DIR, AUDIO_FILENAME)
json_path = os.path.join(OUTPUT_DIR, JSON_FILENAME)


# Pydantic model for transcription response
class TranscriptionResponse(BaseModel):
    timestamp: str
    audio_file: str
    transcription: str


# ==== Document Handlers ====
def extract_text_from_pdf(content):
    reader = PdfReader(io.BytesIO(content))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_text_from_txt(content):
    return content.decode("utf-8")


def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join(para.text for para in doc.paragraphs)


def extract_text_from_pptx(file):
    prs = Presentation(file)
    return "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def extract_text_from_image(file):
    img = Image.open(file)
    return pytesseract.image_to_string(img)


def get_text_from_file(uploaded_file: UploadFile):
    ext = os.path.splitext(uploaded_file.filename)[1].lower()
    content = uploaded_file.file.read()
    uploaded_file.file.seek(0)

    if ext == ".pdf":
        return extract_text_from_pdf(content)
    elif ext == ".txt":
        return extract_text_from_txt(content)
    elif ext == ".docx":
        return extract_text_from_docx(uploaded_file.file)
    elif ext == ".pptx":
        return extract_text_from_pptx(uploaded_file.file)
    elif ext in [".png", ".jpg", ".jpeg"]:
        return extract_text_from_image(uploaded_file.file)
    else:
        return "Unsupported document type"


# ==== Quizzy Bot Logic ====
def generate_response(user_message, history):
    global document_context
    full_query = user_message
    if document_context:
        full_query = f"Based on the following document content:\n\n{document_context}\n\nAnd the user's query: {user_message}"

    try:
        prompt_with_history = ""
        for human, ai in history:
            prompt_with_history += f"Human: {human}\nAI: {ai}\n"
        prompt_with_history += f"Human: {full_query}\nAI:"

        # Sending the prompt to Gemini for processing
        response = model.generate_content(prompt_with_history)
        bot_response = response.text.strip()

        # Format the bot response (add some basic HTML structure for clarity)
        structured_response = format_bot_response(bot_response)

        # Do not append to history here, it's handled in the /chat endpoint
        return structured_response

    except Exception as e:
        return f"An error occurred: {e}. Please try again or rephrase your query."


def format_bot_response(response: str) -> str:
    """
    Formats the Gemini response by:
    - Removing markdown and HTML
    - Ensuring clear line breaks between bullet/numbered points
    - Fixing special characters like &
    - Maintaining clean, readable spacing
    """
    # Remove markdown formatting
    cleaned_response = re.sub(r"\*\*(.*?)\*\*", r"\1", response)
    cleaned_response = re.sub(r"\*(.*?)\*", r"\1", cleaned_response)

    # Remove HTML tags and replace <br> with line breaks
    cleaned_response = cleaned_response.replace("<br>", "\n")
    cleaned_response = re.sub(r"<[^>]+>", "", cleaned_response)

    # Fix encoded characters
    cleaned_response = cleaned_response.replace("\\&", "&").replace("&nbsp;", " ")

    # Add line breaks between numbered or bullet items if not already present
    cleaned_response = re.sub(r"(?<=\d\.)\s*", " ", cleaned_response)  # e.g. 1. Text
    cleaned_response = re.sub(
        r"\n(?=\d\.)", "\n\n", cleaned_response
    )  # Ensure spacing before numbers
    cleaned_response = re.sub(
        r"\n(?=- )", "\n\n", cleaned_response
    )  # Ensure spacing before bullets

    # Normalize multiple newlines
    cleaned_response = re.sub(r"\n{2,}", "\n\n", cleaned_response).strip()

    return cleaned_response


# ==== API Endpoints for Quizzy Bot ====
@app.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    global document_context
    try:
        document_context = get_text_from_file(file)
        # Provide a preview of the uploaded document (first 300 characters)
        return {
            "message": "Document uploaded successfully.",
            "preview": document_context[:300],
        }
    except Exception as e:
        return {"error": f"Error uploading document: {str(e)}"}


@app.post("/chat")
async def chat(message: str = Form(...)):
    global document_context, chat_history
    query = (
        f"Based on the document:\n{document_context}\nUser asked: {message}"
        if document_context
        else message
    )

    try:
        response = generate_response(query, chat_history)
        if not response:
            return {"error": "Empty response from Gemini."}

        chat_history.append(
            (message, response)
        )  # Append to history after successful generation
        return {"response": response}
    except Exception as e:
        return {"error": f"Gemini API error: {str(e)}"}


@app.post("/clear")
async def clear():
    global document_context, chat_history
    chat_history = []
    document_context = ""
    return {"message": "Context and chat history cleared."}


# ==== API Endpoint for Audio Recording ====
@app.post("/record_audio", response_model=TranscriptionResponse)
async def record_audio():
    # Recording the audio
    print(f"\nRecording for {DURATION} seconds... Speak now!")
    audio_data = sd.rec(
        int(DURATION * SAMPLE_RATE), samplerate=SAMPLE_RATE, channels=1, dtype="int16"
    )
    sd.wait()

    # Save the recorded audio
    write(audio_path, SAMPLE_RATE, audio_data)
    print(f"Audio saved to: {audio_path}")

    # Process audio with speech recognition
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        recorded_audio = recognizer.record(source)

    try:
        text = recognizer.recognize_google(recorded_audio)
        print("Transcribed Text:", text)
    except sr.UnknownValueError:
        text = "[Unrecognized speech]"
        print("Could not understand the audio.")
    except sr.RequestError:
        text = "[Google API unavailable]"
        print("Could not request transcription from Google API.")

    # Create output response
    output = {
        "timestamp": datetime.now().isoformat(),
        "audio_file": audio_path,
        "transcription": text,
    }

    # Save transcription to JSON file
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(output, f, indent=4)

    print(f"Transcription saved to: {json_path}")

    # Return the result
    return TranscriptionResponse(**output)

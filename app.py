# from fastapi import FastAPI, File, UploadFile
# from fastapi.middleware.cors import CORSMiddleware  # <-- Import CORSMiddleware here
# from pydantic import BaseModel
# from scipy.io.wavfile import write
# import sounddevice as sd
# import os
# import json
# from datetime import datetime
# import speech_recognition as sr
# import io
# from tempfile import NamedTemporaryFile

# app = FastAPI()

# # CORS configuration
# origins = [
#     "http://localhost:5173",  # React development server
# ]

# # Add CORS middleware
# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=origins,  # Allows requests from React frontend
#     allow_credentials=True,
#     allow_methods=["*"],  # Allows all HTTP methods like GET, POST, etc.
#     allow_headers=["*"],  # Allows all headers
# )

# # Constants
# DURATION = 10  # seconds to record
# SAMPLE_RATE = 44100
# AUDIO_DIR = "audio_files"
# OUTPUT_DIR = "outputs"
# AUDIO_FILENAME = "audio.wav"
# JSON_FILENAME = "transcription.json"

# # Create directories if they don't exist
# os.makedirs(AUDIO_DIR, exist_ok=True)
# os.makedirs(OUTPUT_DIR, exist_ok=True)

# audio_path = os.path.join(AUDIO_DIR, AUDIO_FILENAME)
# json_path = os.path.join(OUTPUT_DIR, JSON_FILENAME)


# class TranscriptionResponse(BaseModel):
#     timestamp: str
#     audio_file: str
#     transcription: str


# @app.post("/record_audio", response_model=TranscriptionResponse)
# async def record_audio():
#     # Recording the audio
#     print(f"\nRecording for {DURATION} seconds... Speak now!")
#     audio_data = sd.rec(
#         int(DURATION * SAMPLE_RATE), samplerate=SAMPLE_RATE, channels=1, dtype="int16"
#     )
#     sd.wait()

#     # Save the recorded audio
#     write(audio_path, SAMPLE_RATE, audio_data)
#     print(f"Audio saved to: {audio_path}")

#     # Process audio with speech recognition
#     recognizer = sr.Recognizer()
#     with sr.AudioFile(audio_path) as source:
#         recorded_audio = recognizer.record(source)

#     try:
#         text = recognizer.recognize_google(recorded_audio)
#         print("Transcribed Text:", text)
#     except sr.UnknownValueError:
#         text = "[Unrecognized speech]"
#         print("Could not understand the audio.")
#     except sr.RequestError:
#         text = "[Google API unavailable]"
#         print("Could not request transcription from Google API.")

#     # Create output response
#     output = {
#         "timestamp": datetime.now().isoformat(),
#         "audio_file": audio_path,
#         "transcription": text,
#     }

#     # Save transcription to JSON file
#     with open(json_path, "w", encoding="utf-8") as f:
#         json.dump(output, f, indent=4)

#     print(f"Transcription saved to: {json_path}")

#     # Return the result
#     return TranscriptionResponse(**output)


import os
import json
import io
import re
from datetime import datetime
from tempfile import NamedTemporaryFile

import sounddevice as sd
import speech_recognition as sr
from scipy.io.wavfile import write

from fastapi import FastAPI, File, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

import google.generativeai as genai
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

# Initialize FastAPI app
app = FastAPI()

# --- CORS Configuration (Combined) ---
origins = [
    "http://localhost:5173",  # React development server
    # Add other origins for production if needed, e.g., your Render frontend URL
    "https://your-render-frontend-url.onrender.com",  # Example for Render deployment
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "*"
    ],  # Allows all origins for simplicity in development, tighten in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- Audio Recording Constants and Setup ---
DURATION = 10  # seconds to record
SAMPLE_RATE = 44100
AUDIO_DIR = "audio_files"
OUTPUT_DIR = "outputs"
AUDIO_FILENAME = "audio.wav"
JSON_FILENAME = "transcription.json"

os.makedirs(AUDIO_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

audio_path = os.path.join(AUDIO_DIR, AUDIO_FILENAME)
json_path = os.path.join(OUTPUT_DIR, JSON_FILENAME)


class TranscriptionResponse(BaseModel):
    timestamp: str
    audio_file: str
    transcription: str


# --- Gemini Configuration ---
# It's highly recommended to use environment variables for API keys in production
# For Render, you can set these in the 'Environment' section of your service.
API_KEY = os.getenv(
    "GEMINI_API_KEY", "AIzaSyA58up6mb0EppG3dI0lT2WYct4Om9aEQKw"
)  # Replace with your actual key or use environment variable
genai.configure(api_key=API_KEY)
model = genai.GenerativeModel(
    "gemini-2.0-flash",
    system_instruction="Your responses should be concise, short, direct, and educational, answering only educational queries with a 120 to 200 word limit.",
)

# --- Global State for Quizzy Bot ---
chat_history = []
document_context = ""


# --- Audio Recording Endpoint ---
@app.post("/record_audio", response_model=TranscriptionResponse)
async def record_audio():
    print(f"\nRecording for {DURATION} seconds... Speak now!")
    audio_data = sd.rec(
        int(DURATION * SAMPLE_RATE), samplerate=SAMPLE_RATE, channels=1, dtype="int16"
    )
    sd.wait()

    write(audio_path, SAMPLE_RATE, audio_data)
    print(f"Audio saved to: {audio_path}")

    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        recorded_audio = recognizer.record(source)

    try:
        text = recognizer.recognize_google(recorded_audio)
        print("Transcribed Text:", text)
    except sr.UnknownValueError:
        text = "[Unrecognized speech]"
        print("Could not understand the audio.")
    except sr.RequestError:
        text = "[Google API unavailable]"
        print("Could not request transcription from Google API.")

    output = {
        "timestamp": datetime.now().isoformat(),
        "audio_file": audio_path,
        "transcription": text,
    }

    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(output, f, indent=4)

    print(f"Transcription saved to: {json_path}")

    return TranscriptionResponse(**output)


# --- Document Handlers ---
def extract_text_from_pdf(content):
    reader = PdfReader(io.BytesIO(content))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_text_from_txt(content):
    return content.decode("utf-8")


def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join(para.text for para in doc.paragraphs)


def extract_text_from_pptx(file):
    prs = Presentation(file)
    return "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def extract_text_from_image(file):
    img = Image.open(file)
    # Ensure pytesseract is configured if not in PATH. For Render, ensure it's installed via buildpack or similar.
    # pytesseract.pytesseract.tesseract_cmd = r'/usr/bin/tesseract' # Example path for Linux
    return pytesseract.image_to_string(img)


def get_text_from_file(uploaded_file: UploadFile):
    ext = os.path.splitext(uploaded_file.filename)[1].lower()
    content = uploaded_file.file.read()
    uploaded_file.file.seek(0)  # Reset file pointer for subsequent reads if any

    if ext == ".pdf":
        return extract_text_from_pdf(content)
    elif ext == ".txt":
        return extract_text_from_txt(content)
    elif ext == ".docx":
        return extract_text_from_docx(io.BytesIO(content))  # Pass BytesIO for docx
    elif ext == ".pptx":
        return extract_text_from_pptx(io.BytesIO(content))  # Pass BytesIO for pptx
    elif ext in [".png", ".jpg", ".jpeg"]:
        return extract_text_from_image(io.BytesIO(content))  # Pass BytesIO for image
    else:
        return "Unsupported document type"


# --- Quizzy Bot Logic ---
def generate_response(user_message, history):
    global document_context
    full_query = user_message
    if document_context:
        full_query = f"Based on the following document content:\n\n{document_context}\n\nAnd the user's query: {user_message}"

    try:
        prompt_with_history = ""
        for human, ai in history:
            prompt_with_history += f"Human: {human}\nAI: {ai}\n"
        prompt_with_history += f"Human: {full_query}\nAI:"

        response = model.generate_content(prompt_with_history)
        bot_response = response.text.strip()

        structured_response = format_bot_response(bot_response)

        # Update history *after* getting the response, not inside the function that generates it
        # history.append((user_message, structured_response)) # This should be handled by the caller (chat endpoint)
        return structured_response

    except Exception as e:
        print(f"Error generating response: {e}")  # Log the error
        return f"An error occurred: {e}. Please try again or rephrase your query."


def format_bot_response(response: str) -> str:
    """
    Formats the Gemini response by:
    - Removing markdown and HTML
    - Ensuring clear line breaks between bullet/numbered points
    - Fixing special characters like \&
    - Maintaining clean, readable spacing
    """
    cleaned_response = re.sub(r"\*\*(.*?)\*\*", r"\1", response)
    cleaned_response = re.sub(r"\*(.*?)\*", r"\1", cleaned_response)

    cleaned_response = cleaned_response.replace("<br>", "\n")
    cleaned_response = re.sub(r"<[^>]+>", "", cleaned_response)

    cleaned_response = cleaned_response.replace("\\&", "&").replace("&nbsp;", " ")

    cleaned_response = re.sub(r"(?<=\d\.)\s*", " ", cleaned_response)
    cleaned_response = re.sub(r"\n(?=\d\.)", "\n\n", cleaned_response)
    cleaned_response = re.sub(r"\n(?=- )", "\n\n", cleaned_response)

    cleaned_response = re.sub(r"\n{2,}", "\n\n", cleaned_response).strip()

    return cleaned_response


# --- Quizzy Bot API Endpoints ---
@app.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    global document_context
    try:
        document_context = get_text_from_file(file)
        if document_context == "Unsupported document type":
            return {
                "error": "Unsupported document type. Please upload PDF, TXT, DOCX, PPTX, PNG, JPG, or JPEG files."
            }

        return {
            "message": "Document uploaded successfully.",
            "preview": (
                document_context[:300] + "..."
                if len(document_context) > 300
                else document_context
            ),
        }
    except Exception as e:
        print(f"Error during document upload: {e}")  # Log the error
        return {"error": f"Error uploading document: {str(e)}"}


@app.post("/chat")
async def chat(message: str = Form(...)):
    global document_context, chat_history
    query = (
        f"Based on the document:\n{document_context}\nUser asked: {message}"
        if document_context
        else message
    )

    try:
        response = generate_response(query, chat_history)
        if not response:
            return {"error": "Empty response from Gemini."}

        # Append to chat history after a successful response
        chat_history.append((message, response))
        return {"response": response}
    except Exception as e:
        print(f"Error in chat endpoint: {e}")  # Log the error
        return {"error": f"Gemini API error: {str(e)}"}


@app.post("/clear")
async def clear():
    global document_context, chat_history
    chat_history = []
    document_context = ""
    return {"message": "Context and chat history cleared."}


# To run this with Uvicorn (for Render deployment):
# 1. Save this file as, e.g., `main.py`.
# 2. Make sure you have the necessary packages installed:
#    `pip install fastapi uvicorn python-multipart python-dotenv pydantic scipy sounddevice SpeechRecognition PyPDF2 python-docx python-pptx pillow pytesseract google-generativeai`
#    You might also need `portaudio` for `sounddevice` depending on your OS.
#    For `pytesseract`, you'll need the Tesseract OCR engine installed on your system. On Render, you might need a custom buildpack or a Dockerfile to include it.
# 3. Create a `requirements.txt` file with all dependencies.
# 4. In your Render service, set the "Start Command" to `uvicorn main:app --host 0.0.0.0 --port $PORT`.
# 5. Set your `GEMINI_API_KEY` as an environment variable in Render.
